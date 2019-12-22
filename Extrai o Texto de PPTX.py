# pip3.6 install python-pptx

from pptx import Presentation

prs = Presentation("nomeDoArquivo.pptx")

for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print(shape.text)
